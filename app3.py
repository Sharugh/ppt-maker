import streamlit as st
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import requests
from io import BytesIO
from PIL import Image

# Function to extract text from a Word document
def extract_text_from_docx(file):
    doc = Document(file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return "\n".join(full_text)

# Function to search for images using Pexels API
def search_image(query):
    api_key = "LnDqCT2BbahSLy0YduB4sfFyVC0P8EBnJREHaxsXBUs"  # Replace with your Pexels API key
    url = f"https://api.pexels.com/v1/search?query={query}&per_page=1"
    headers = {"Authorization": api_key}
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        results = response.json()
        if results["photos"]:
            return results["photos"][0]["src"]["large"]  # Return the first image URL
    return None

# Function to create a professional PPT with images
def create_ppt(text, output_file):
    prs = Presentation()

    # Use a professional template (optional)
    # prs = Presentation("professional_template.pptx")

    # Slide layout for title and content
    slide_layout = prs.slide_layouts[1]

    # Split text into slides (e.g., by sections)
    slides_content = text.split("\n\n")

    # Add a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "India Energy Week 2025"
    subtitle.text = "Key Discussions and Outcomes"

    # Add content slides
    for content in slides_content:
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]

        # Set title and body text
        title.text = "Slide Title"  # Customize based on content
        body.text = content

        # Extract keywords for image search
        keywords = content.split()[:5]  # Use the first 5 words as keywords
        query = " ".join(keywords)
        image_url = search_image(query)

        # Add image to the slide if found
        if image_url:
            try:
                img = download_image(image_url)
                slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(4))
            except Exception as e:
                st.warning(f"Could not add image for query: {query}. Error: {e}")

        # Add hyperlinks for sources
        if "Source:" in content:
            source_text = content.split("Source:")[1].strip()
            source_url = source_text.split(" ")[0]  # Extract URL
            run = body.text_frame.add_paragraph().add_run()
            run.text = "Source"
            run.hyperlink.address = source_url

    # Save the PPT
    prs.save(output_file)

# Function to download an image from a URL
def download_image(url):
    response = requests.get(url)
    return BytesIO(response.content)

# Streamlit app
st.title("Word to Professional PPT Converter")
st.write("Upload a Word document to generate a visually appealing PPT.")

# File upload
uploaded_file = st.file_uploader("Upload a Word document", type=["docx"])

if uploaded_file is not None:
    # Extract text from the Word document
    text = extract_text_from_docx(uploaded_file)

    # Create a PPT
    ppt_file = "output_presentation.pptx"
    create_ppt(text, ppt_file)

    # Provide download links
    st.success("Professional PPT generated successfully!")
    with open(ppt_file, "rb") as f:
        st.download_button("Download PPT", f, file_name=ppt_file)
